from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE

ROOT=Path('/home/user/mas_reliab'); FIG=ROOT/'figures'
NAVY='0B1F33'; BLUE='146C94'; TEAL='1B998B'; CYAN='54C6EB'; AMBER='F4A261'; RED='D95D5D'; GREEN='2A9D8F'; LIGHT='F3F6F8'; MID='DCE5EA'; DARK='23323D'; GRAY='657681'; PURPLE='7A5AF8'; WHITE='FFFFFF'

def rgb(h): return RGBColor.from_string(h)
prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
blank=prs.slide_layouts[6]

def rect(sl,x,y,w,h,fill=WHITE,line=None,radius=False):
    sh=sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=rgb(fill)
    sh.line.color.rgb=rgb(line or fill)
    if radius: sh.adjustments[0]=0.08
    return sh

def textbox(sl,x,y,w,h,text,size=20,color=DARK,bold=False,font='Aptos',align=PP_ALIGN.LEFT,val=MSO_ANCHOR.TOP,margin=.06):
    tb=sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)); tf=tb.text_frame; tf.clear(); tf.word_wrap=True; tf.vertical_anchor=val
    tf.margin_left=tf.margin_right=tf.margin_top=tf.margin_bottom=Inches(margin)
    p=tf.paragraphs[0]; p.text=text; p.alignment=align; p.font.name=font; p.font.size=Pt(size); p.font.bold=bold; p.font.color.rgb=rgb(color)
    return tb

def rich_text(sl,x,y,w,h,runs,size=20,align=PP_ALIGN.LEFT,val=MSO_ANCHOR.TOP):
    tb=sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)); tf=tb.text_frame; tf.clear(); tf.word_wrap=True; tf.vertical_anchor=val
    tf.margin_left=tf.margin_right=Inches(.04)
    p=tf.paragraphs[0]; p.alignment=align
    for text,color,bold in runs:
        r=p.add_run(); r.text=text; r.font.name='Aptos'; r.font.size=Pt(size); r.font.color.rgb=rgb(color); r.font.bold=bold
    return tb

def title(sl,t,sub=None,num=None):
    textbox(sl,.62,.30,11.9,.55,t,25,NAVY,True)
    rect(sl,.63,.91,1.05,.055,TEAL)
    if sub: textbox(sl,.63,1.03,11.7,.38,sub,11.5,GRAY,False)
    if num is not None: textbox(sl,12.42,.34,.35,.3,str(num),9,GRAY,False,align=PP_ALIGN.RIGHT)

def footer(sl,source=None,status=False):
    rect(sl,.62,7.15,12.05,.012,MID)
    textbox(sl,.63,7.20,6.1,.18,'MAS-RELIAB  •  Dev Parth  •  Viva Defence',7.5,GRAY)
    if source: textbox(sl,6.5,7.20,6.17,.18,source,7.2,GRAY,align=PP_ALIGN.RIGHT)
    if status:
        sh=rect(sl,10.28,6.72,2.38,.34,'FDECEC',RED,True); textbox(sl,10.36,6.77,2.22,.19,'NO EMPIRICAL RESULT YET',8.3,RED,True,align=PP_ALIGN.CENTER)

def bullets(sl,items,x=.83,y=1.55,w=11.7,h=5.15,size=19,color=DARK,accent=TEAL,gap=9):
    tb=sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)); tf=tb.text_frame; tf.clear(); tf.word_wrap=True
    tf.margin_left=tf.margin_right=Inches(.02)
    for i,item in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.level=0; p.space_after=Pt(gap); p.text=''
        r=p.add_run(); r.text='● '; r.font.color.rgb=rgb(accent); r.font.size=Pt(size-2); r.font.name='Aptos'
        if isinstance(item,tuple):
            r=p.add_run(); r.text=item[0]; r.font.bold=True; r.font.color.rgb=rgb(color); r.font.size=Pt(size); r.font.name='Aptos'
            r=p.add_run(); r.text=item[1]; r.font.color.rgb=rgb(color); r.font.size=Pt(size); r.font.name='Aptos'
        else:
            r=p.add_run(); r.text=item; r.font.color.rgb=rgb(color); r.font.size=Pt(size); r.font.name='Aptos'
    return tb

def pill(sl,x,y,w,text,color=TEAL):
    rect(sl,x,y,w,.42,color,color,True); textbox(sl,x+.03,y+.06,w-.06,.22,text,10.5,WHITE,True,align=PP_ALIGN.CENTER)

def card(sl,x,y,w,h,head,body,color=TEAL,head_size=14,body_size=11.5):
    rect(sl,x,y,w,h,WHITE,MID,True); rect(sl,x,y,.07,h,color,color,True)
    textbox(sl,x+.22,y+.16,w-.35,.35,head,head_size,color,True)
    textbox(sl,x+.22,y+.60,w-.35,h-.72,body,body_size,DARK)

def image(sl,file,x,y,w,h=None):
    if h: sl.shapes.add_picture(str(FIG/file),Inches(x),Inches(y),Inches(w),Inches(h))
    else: sl.shapes.add_picture(str(FIG/file),Inches(x),Inches(y),width=Inches(w))

def new_slide():
    sl=prs.slides.add_slide(blank); rect(sl,0,0,13.333,7.5,WHITE); return sl

# 1 title
sl=new_slide(); rect(sl,0,0,13.333,7.5,NAVY)
rect(sl,.7,.7,.12,5.9,TEAL)
textbox(sl,1.18,1.05,10.9,1.65,'A Framework for Evaluating\nReliability and Failure Propagation\nin Multi-Agent AI Systems',29,WHITE,True)
textbox(sl,1.20,3.15,5.3,.65,'MAS-RELIAB',25,TEAL,True)
textbox(sl,1.20,3.83,8.8,.55,'Research Paper and Viva Defence',17,'DCE5EA')
textbox(sl,1.20,4.55,5.0,.48,'Dev Parth',20,WHITE,True)
textbox(sl,1.20,5.06,8.0,.38,'Pre-experimental, results-ready protocol  •  16 August 2026',11.5,'B9C7CE')
pill(sl,10.08,5.76,2.30,'DEFENCE DECK',TEAL)

# 2 integrity
sl=new_slide(); title(sl,'Scientific status','What this defence does—and does not—claim',2)
card(sl,.75,1.55,5.75,3.65,'Complete now','Formal framework • task schema • fault operators • metrics • six experiments • statistical plan • results shell • viva defence',TEAL,18,16)
card(sl,6.83,1.55,5.75,3.65,'Pending empirical execution','No run traces • no estimates • no confidence intervals • no p-values • no hypothesis outcomes',RED,18,16)
rect(sl,.76,5.55,11.82,.72,'FFF3E7',AMBER,True); textbox(sl,1.00,5.76,11.35,.30,'Every Results value remains “—” until populated from the frozen analysis pipeline.',16,RED,True,align=PP_ALIGN.CENTER)
footer(sl,status=True)

# 3 problem
sl=new_slide(); title(sl,'A local error can become a system failure','Multi-agent collaboration creates dependency paths',3)
image(sl,'figure_02_graph.png',.70,1.35,7.7)
card(sl,8.72,1.60,3.85,1.15,'1. Origin','A wrong fact, truncated handoff, stale tool result, or manager decision.',RED,13,11)
card(sl,8.72,2.95,3.85,1.15,'2. Transmission','Downstream agents consume and reuse the altered state.',AMBER,13,11)
card(sl,8.72,4.30,3.85,1.15,'3. Outcome','The team may converge on a coherent but invalid end state.',PURPLE,13,11)
textbox(sl,8.78,5.72,3.75,.56,'Question: where did it start, how far did it travel, and what contained it?',13,NAVY,True,align=PP_ALIGN.CENTER)
footer(sl)

# 4 why success not enough
sl=new_slide(); title(sl,'Why final task success is insufficient','A binary score collapses operationally different failures',4)
items=[('Repeatability','Passed once; does it pass across k repeats?'),('Robustness','What happens under paraphrase, tool faults, or stale state?'),('Propagation','Did one fault spread to one node or the whole team?'),('Attribution','Can we find the responsible agent and decisive step?'),('Mitigation','Did verification or recovery actually contain the lineage?'),('Resources','How many tokens, calls, seconds, and currency units?')]
for i,(h,b) in enumerate(items):
    x=.78+(i%3)*4.08; y=1.50+(i//3)*2.13; card(sl,x,y,3.65,1.72,h,b,[TEAL,BLUE,PURPLE,RED,AMBER,GRAY][i],16,12.5)
footer(sl)

# 5 related work
sl=new_slide(); title(sl,'Related work covers major components','MAS-RELIAB must be positioned as integration, not priority',5)
works=[('MAST','failure taxonomy'),('τ-bench','state + pass^k'),('ReliabilityBench','repeats + tool faults'),('TraceElephant','full-trace attribution'),('AgenTracer / ErrorProbe','origin diagnosis'),('AgentAsk','handoff containment'),('RiskLab / TAMAS','controlled risk'),('MultiAgentBench / SILO','topology + coordination')]
for i,(h,b) in enumerate(works):
    x=.75+(i%4)*3.10; y=1.46+(i//4)*1.48; card(sl,x,y,2.83,1.16,h,b,[TEAL,BLUE,PURPLE,AMBER,RED,GREEN,CYAN,GRAY][i],12.5,10.5)
rect(sl,.75,4.65,11.92,1.28,'FDECEC',RED,True)
textbox(sl,1.02,4.87,2.55,.34,'Closest overlap',14,RED,True)
textbox(sl,3.06,4.82,9.23,.50,'From Spark to Fire already uses directed dependency graphs, atomic error seeds, cascade analysis, topology sensitivity, and mitigation.',15,DARK,True)
textbox(sl,1.03,6.12,11.3,.37,'Therefore: no claim that graph-based propagation or cascade mitigation is novel.',13,RED,True,align=PP_ALIGN.CENTER)
footer(sl,'[1–15] in manuscript')

# 6 novelty
sl=new_slide(); title(sl,'Defensible novelty boundary','The contribution is the controlled linkage of dimensions',6)
image(sl,'figure_01_architecture.png',.65,1.35,7.15)
card(sl,8.08,1.47,4.52,3.83,'MAS-RELIAB contribution','One paired protocol connects:\n\n• multidimensional reliability\n• verified fault injection\n• graph-normalized propagation\n• evidence-view attribution\n• verification and recovery\n• budget-matched topology\n• tokens, latency, and cost',TEAL,17,14)
rect(sl,8.08,5.56,4.52,.64,'FFF3E7',AMBER,True); textbox(sl,8.25,5.73,4.18,.25,'Methodology claim ≠ empirical superiority',12,RED,True,align=PP_ALIGN.CENTER)
footer(sl)

# 7 aim objectives
sl=new_slide(); title(sl,'Aim and objectives','Turn “reliable” into reproducible measurements and decisions',7)
textbox(sl,.83,1.35,11.7,.72,'Aim: design and empirically validate MAS-RELIAB as a reproducible method for measuring reliability and failure propagation in multi-agent AI systems.',18,NAVY,True,align=PP_ALIGN.CENTER,val=MSO_ANCHOR.MIDDLE)
bullets(sl,[
 ('Represent ','agents, events, tools, and state as time-aware dependencies.'),
 ('Inject ','documented faults with applied and consumed checks.'),
 ('Measure ','outcome, repeatability, robustness, propagation, attribution, recovery, and resources.'),
 ('Compare ','single, sequential, parallel, and hierarchical architectures under controls.'),
 ('Release ','schemas, manifests, analysis rules, and results-ready artifacts.')],x=1.05,y=2.40,w=11.0,h=3.9,size=17,gap=10)
footer(sl)

# 8 RQ grid
sl=new_slide(); title(sl,'Research questions RQ1–RQ6','Each question maps to one confirmatory experiment',8)
rqs=[('RQ1','Reliability','Beyond one success score'),('RQ2','Propagation','Fault type, severity, position'),('RQ3','Attribution','Effect of observability'),('RQ4','Verification','Local vs final-only'),('RQ5','Recovery','Reliability–cost trade-off'),('RQ6','Topology','Architecture × fault location')]
for i,(n,h,b) in enumerate(rqs):
    x=.80+(i%3)*4.12; y=1.50+(i//3)*2.15
    rect(sl,x,y,3.70,1.74,WHITE,MID,True); pill(sl,x+.18,y+.18,.72,n,[TEAL,BLUE,PURPLE,RED,AMBER,GRAY][i]); textbox(sl,x+1.05,y+.19,2.35,.30,h,15,NAVY,True); textbox(sl,x+.22,y+.82,3.20,.52,b,12,DARK,False,align=PP_ALIGN.CENTER)
footer(sl)

# 9 hypotheses
sl=new_slide(); title(sl,'Preregistered hypotheses H1–H6','Falsifiable claims—not expected results',9)
hs=[('H1','Earlier faults → greater normalized impact/depth'),('H2','Full traces → better agent + step attribution'),('H3','Local verification → lower EPR and DAF'),('H4','Rollback / alternate agent → more recovery, more cost'),('H5','Topology interacts with fault location'),('H6','Multidimensional profile reveals masked weaknesses')]
for i,(h,b) in enumerate(hs): card(sl,.83+(i%2)*6.05,1.38+(i//2)*1.62,5.62,1.28,h,b,[RED,TEAL,BLUE,AMBER,PURPLE,GRAY][i],14,12.5)
textbox(sl,.95,6.38,11.45,.34,'Null, mixed, or contradictory findings will be reported without reframing.',13,RED,True,align=PP_ALIGN.CENTER)
footer(sl,status=True)

# 10 architecture
sl=new_slide(); title(sl,'Separation of concerns','AI Command Center executes; MAS-RELIAB defines the science',10)
image(sl,'figure_01_architecture.png',1.05,1.28,11.25)
footer(sl)

# 11 formal episode
sl=new_slide(); title(sl,'Formal episode model','One immutable record links configuration, intervention, trace, outcome, and cost',11)
textbox(sl,.80,1.35,11.75,.60,'𝓔 = ⟨ τ, G, Π, Ω, I, X, Y, C ⟩',25,NAVY,True,font='Cambria Math',align=PP_ALIGN.CENTER)
labels=[('τ','versioned task'),('G','dependency graph'),('Π','model / roles / policy'),('Ω','environment snapshot'),('I','fault or null intervention'),('X','ordered trace'),('Y','validated outcome'),('C','resource record')]
for i,(s,b) in enumerate(labels):
    x=.78+(i%4)*3.11; y=2.25+(i//4)*1.40
    rect(sl,x,y,2.80,1.05,WHITE,MID,True); textbox(sl,x+.15,y+.16,.50,.36,s,19,TEAL,True,font='Cambria Math',align=PP_ALIGN.CENTER); textbox(sl,x+.70,y+.19,1.92,.50,b,11.5,DARK,True,align=PP_ALIGN.CENTER,val=MSO_ANCHOR.MIDDLE)
rect(sl,1.35,5.39,10.65,.63,'EAF5F3',TEAL,True); textbox(sl,1.55,5.56,10.25,.26,'Valid injection = applied as specified AND consumed by the intended component.',14,TEAL,True,align=PP_ALIGN.CENTER)
footer(sl)

# 12 graph
sl=new_slide(); title(sl,'Propagation is measured on reachable dependencies','Temporal succession alone is not causal evidence',12)
image(sl,'figure_02_graph.png',.72,1.27,7.75)
card(sl,8.72,1.48,3.88,1.18,'EPR','Fraction of valid injections causing any downstream failure.',RED,14,11.5)
card(sl,8.72,2.89,3.88,1.18,'DAF','Affected downstream nodes / reachable downstream nodes.',AMBER,14,11.5)
card(sl,8.72,4.30,3.88,1.18,'Depth + amplification','How far and how broadly the lineage spreads.',PURPLE,14,11.5)
textbox(sl,8.75,5.78,3.80,.43,'Primary evidence: instrumented consumption + lineage + task-relevant error.',11.5,NAVY,True,align=PP_ALIGN.CENTER)
footer(sl)

# 13 taxonomy
sl=new_slide(); title(sl,'Operational fault taxonomy','Category names organize analysis; operator definitions create reproducibility',13)
image(sl,'figure_03_taxonomy.png',.75,1.30,11.85)
footer(sl,'Related to MAST and AgentAsk')

# 14 topologies
sl=new_slide(); title(sl,'Architecture baselines','Primary analysis isolates topology with matched budgets',14)
image(sl,'figure_04_topologies.png',.70,1.42,11.95)
rect(sl,1.08,6.15,11.18,.54,'EAF5F3',TEAL,True); textbox(sl,1.25,6.29,10.84,.24,'Same tasks • seed blocks • model family • information • validator • token/tool/timeout caps',12.5,TEAL,True,align=PP_ALIGN.CENTER)
footer(sl)

# 15 tasks
sl=new_slide(); title(sl,'Benchmark tasks and ground truth','Objectively checkable end states are the primary standard',15)
tasks=[('Research synthesis','Closed corpus + fact/evidence keys',TEAL),('Data analysis','Reference computation + tolerances',BLUE),('Software debugging','Hidden/public executable tests',PURPLE),('Decision support','Constraint solver + feasibility',AMBER),('Stateful tool workflow','Database end-state equivalence',RED)]
for i,(h,b,c) in enumerate(tasks):
    x=.77+i*2.50; rect(sl,x,1.55,2.25,3.35,WHITE,MID,True); rect(sl,x,1.55,2.25,.68,c,c,True); textbox(sl,x+.10,1.74,2.05,.24,h,12,WHITE,True,align=PP_ALIGN.CENTER); textbox(sl,x+.20,2.68,1.85,1.10,b,13,DARK,True,align=PP_ALIGN.CENTER,val=MSO_ANCHOR.MIDDLE)
textbox(sl,.94,5.35,11.48,.66,'LLM-as-a-judge may be a secondary diagnostic; it is not the sole ground truth for the core benchmark.',14,RED,True,align=PP_ALIGN.CENTER,val=MSO_ANCHOR.MIDDLE)
footer(sl)

# 16 metrics
sl=new_slide(); title(sl,'Multidimensional reliability vector','Report components first; use Pareto trade-offs—not hidden weights',16)
image(sl,'figure_06_metrics.png',.70,1.28,11.90)
footer(sl)

# 17 formulas
sl=new_slide(); title(sl,'Core propagation and reliability metrics','Explicit denominators prevent misleading claims',17)
forms=[('TSR','N successful / N eligible episodes'),('passᵏ','Tasks whose k prespecified repeats all pass'),('EPR','Episodes with any downstream effect / valid consumed injections'),('DAF','Affected reachable nodes / reachable downstream nodes'),('Recovery','Correct end state after recovery / valid consumed faults'),('Attribution','Agent exact • step exact • joint • top-k • distance')]
for i,(h,b) in enumerate(forms): card(sl,.80+(i%2)*6.05,1.40+(i//2)*1.65,5.62,1.28,h,b,[TEAL,BLUE,RED,AMBER,PURPLE,GRAY][i],15,12)
textbox(sl,1.05,6.43,11.25,.28,'For one serial path: traversal ≈ product of edge probabilities; multi-path independence is not assumed in primary evidence.',11,GRAY,False,align=PP_ALIGN.CENTER)
footer(sl)

# 18 observability
sl=new_slide(); title(sl,'Attribution is conditioned on available evidence','Same failed episodes, three paired trace views',18)
image(sl,'figure_07_observability.png',.70,1.27,11.93)
footer(sl,'TraceElephant and AgenTracer motivate explicit evidence conditions')

# 19 experiment program
sl=new_slide(); title(sl,'Six-experiment program','Every experiment maps to one RQ and hypothesis',19)
image(sl,'figure_05_experiments.png',.73,1.28,11.88)
footer(sl)

# 20 E1/E2
sl=new_slide(); title(sl,'Experiments 1–2','Baseline profile and controlled fault propagation',20)
card(sl,.78,1.48,5.75,4.45,'E1 • Baseline reliability','Four architectures × locked tasks × repeated seeds\n\nPrimary: TSR, pass^k, state consistency\nSecondary: constraint success, tokens, tools, latency, cost\n\nH6: test preregistered ranking discordance / masked weakness',TEAL,18,14)
card(sl,6.80,1.48,5.75,4.45,'E2 • Fault type and position','Operator × severity × early/middle/late depth\n\nPrimary: EPR, DAF, depth, Δ task success\nControls: same task/seed/snapshot; reachable-set size\n\nH1: early vs late after normalization',RED,18,14)
footer(sl)

# 21 E3/E4
sl=new_slide(); title(sl,'Experiments 3–4','Observability for diagnosis; placement for containment',21)
card(sl,.78,1.48,5.75,4.45,'E3 • Attribution','Output-only vs partial vs full trace\n\nSame failures + same gold origins\nPrimary: joint agent/step accuracy, localization distance\nSecondary: top-k, MRR, abstention, diagnostic cost\n\nH2: evidence-view effect',PURPLE,18,14)
card(sl,6.80,1.48,5.75,4.45,'E4 • Verification','None vs final-only vs local-before-handoff\n\nPrimary: EPR, DAF\nSecondary: TSR, false accept/reject, tokens, latency\nNo privileged gold answers\n\nH3: local containment',BLUE,18,14)
footer(sl)

# 22 E5/E6
sl=new_slide(); title(sl,'Experiments 5–6','Recovery trade-offs and topology sensitivity',22)
card(sl,.78,1.48,5.75,4.45,'E5 • Recovery','None • same-context retry • error signal • rollback • alternate agent • human gate\n\nPrimary: recovery rate + residual DAF\nSecondary: repeated-error similarity + resources\nDecision artifact: Pareto frontier\n\nH4: recovery gain with added cost',AMBER,18,13.5)
card(sl,6.80,1.48,5.75,4.45,'E6 • Topology','Single • sequential • parallel • hierarchical\n\nCompare worker vs aggregator and leaf vs manager faults\nPrimary: topology × functional-location interaction\nBudget-matched primary; natural-budget secondary\n\nH5: location-dependent topology effect',TEAL,18,13.5)
footer(sl)

# 23 stats
sl=new_slide(); title(sl,'Statistical analysis plan','Task-level inference, pairing, uncertainty, and effect size',23)
bullets(sl,[
 ('Unit of inference: ','task instance; repeats and faults are nested.'),
 ('Pairing: ','same task, seed block, environment, and fault where meaningful.'),
 ('Tests: ','exact McNemar; paired permutation/Wilcoxon; mixed models.'),
 ('Uncertainty: ','task-clustered 95% confidence intervals.'),
 ('Multiplicity: ','Holm correction within each RQ family.'),
 ('Magnitude: ','risk difference, matched odds ratio, median shift, Cliff’s delta.'),
 ('Power: ','pilot simulation; 100 paired IDs is a floor, not a guarantee.')],x=.95,y=1.35,w=11.4,h=5.45,size=16.3,gap=5)
footer(sl)

# 24 results status
sl=new_slide(); title(sl,'Results section is intentionally locked','No values are shown because no executions were supplied',24)
rect(sl,.82,1.46,11.70,4.28,'F7F9FA',MID,True)
headers=['Architecture','TSR','passᵏ','DAF','Attribution','Cost']
for i,h in enumerate(headers):
    x=1.05+i*1.88; rect(sl,x,1.82,1.72,.52,NAVY,NAVY,True); textbox(sl,x+.05,1.97,1.62,.20,h,10.5,WHITE,True,align=PP_ALIGN.CENTER)
for r in range(4):
    for i in range(6):
        x=1.05+i*1.88; y=2.45+r*.62; rect(sl,x,y,1.72,.50,WHITE,MID,False); textbox(sl,x+.05,y+.14,1.62,.18,'—' if i else ['Single','Sequential','Parallel','Hierarchy'][r],10,DARK,i==0,align=PP_ALIGN.CENTER)
textbox(sl,1.02,5.95,11.28,.46,'Populate only from signed manifests and frozen analysis outputs. Rewrite the abstract, discussion, and conclusion after validation.',14,RED,True,align=PP_ALIGN.CENTER)
footer(sl,status=True)

# 25 reproducibility
sl=new_slide(); title(sl,'Reproducibility by design','Every run must be reconstructable—or its limits stated',25)
cols=[('Configuration','model/version • prompts • topology • budgets • seeds',TEAL),('Environment','snapshot • tools • container • region • hardware',BLUE),('Intervention','operator • target • dose • patch • applied/consumed checks',RED),('Evidence','trace • state deltas • graph • validators • hashes',PURPLE),('Resources','tokens • calls • latency • cost • pricing date',AMBER),('Analysis','manifest • lockfile • code hash • CIs • deviations',GRAY)]
for i,(h,b,c) in enumerate(cols): card(sl,.78+(i%3)*4.10,1.48+(i//3)*2.16,3.68,1.75,h,b,c,15,11.8)
footer(sl)

# 26 limitations
sl=new_slide(); title(sl,'Limitations and threats to validity','The protocol reduces uncertainty; it does not remove it',26)
limits=[('Fault realism','Injected faults may not match natural failure distributions.'),('Graph completeness','Latent dependencies and semantic edges remain uncertain.'),('Counterfactual replay','Stochastic trajectories can change after a repair.'),('Provider drift','Hosted models, policies, load, and prices can change.'),('Observability','Privacy and proprietary constraints may limit traces.'),('Generalization','Selected tasks/topologies cannot represent every deployment.')]
for i,(h,b) in enumerate(limits): card(sl,.82+(i%2)*6.04,1.38+(i//2)*1.65,5.60,1.28,h,b,[RED,TEAL,PURPLE,AMBER,BLUE,GRAY][i],14,12)
footer(sl)

# 27 conclusion
sl=new_slide(); title(sl,'Conclusion','What is defensible today—and what remains to be done',27)
card(sl,.85,1.47,5.65,3.72,'Defensible today','MAS-RELIAB is a complete, auditable, results-ready protocol connecting reliability, verified faults, propagation, attribution, mitigation, topology, and cost.',TEAL,18,15)
card(sl,6.84,1.47,5.65,3.72,'Not yet defensible','No claim that H1–H6 are supported; no superiority claim; no priority claim for graph-based cascades or mitigation.',RED,18,15)
rect(sl,.85,5.53,11.64,.76,NAVY,NAVY,True); textbox(sl,1.15,5.75,11.03,.30,'Next step: pilot → preregistration lock → interleaved execution → validation → populated paper.',16,WHITE,True,align=PP_ALIGN.CENTER)
footer(sl,status=True)

# 28 Q
sl=new_slide(); rect(sl,0,0,13.333,7.5,NAVY)
textbox(sl,1.20,1.48,10.93,.82,'Questions',37,WHITE,True,align=PP_ALIGN.CENTER)
textbox(sl,1.30,2.58,10.73,.70,'MAS-RELIAB',23,TEAL,True,align=PP_ALIGN.CENTER)
textbox(sl,2.06,3.52,9.20,1.08,'A unified experimental methodology for failure propagation, attribution, mitigation, topology, and reliability–cost trade-offs.',17,'DCE5EA',False,align=PP_ALIGN.CENTER)
textbox(sl,4.48,5.15,4.38,.45,'Dev Parth',17,WHITE,True,align=PP_ALIGN.CENTER)

# appendix 29
sl=new_slide(); title(sl,'Appendix • Metric definitions','Quick reference for examiner questions',29)
metric_text=[('TSR','eligible successes / eligible episodes'),('passᵏ','all k repeats pass for a task'),('Consistency','pairwise end-state equivalence'),('Δfault','TSR clean − TSR fault'),('EPR','any downstream effect / valid injections'),('DAF','affected / reachable downstream'),('AF','affected / initially corrupted nodes'),('Depth','max shortest-path distance'),('Recovery','correct after recovery / consumed faults'),('Attribution','agent + step exact/top-k/distance'),('FCS','1 − DAF'),('Resources','tokens, calls, latency, storage, cost')]
for i,(h,b) in enumerate(metric_text):
    x=.73+(i%3)*4.18; y=1.30+(i//3)*1.28; card(sl,x,y,3.78,1.02,h,b,[TEAL,BLUE,PURPLE,RED,AMBER,GRAY][i%6],12,10)
footer(sl)

# appendix 30
sl=new_slide(); title(sl,'Appendix • Hypothesis decision matrix','Primary metrics and contrasts are frozen before unblinding',30)
rows=[('H1','DAF + depth','early vs late'),('H2','joint attribution','full vs partial/output'),('H3','EPR + DAF','local vs final/none'),('H4','recovery + cost','rollback/alternate vs retry'),('H5','interaction','topology × location'),('H6','rank discordance','TSR vs reliability profile')]
for i,(h,m,c) in enumerate(rows):
    y=1.35+i*.83; pill(sl,.80,y,.70,h,[RED,TEAL,BLUE,AMBER,PURPLE,GRAY][i]); textbox(sl,1.75,y+.02,3.00,.32,m,13,NAVY,True); textbox(sl,4.95,y+.02,6.82,.32,c,13,DARK,False); rect(sl,.80,y+.54,11.72,.02,MID)
footer(sl)

# appendix 31
sl=new_slide(); title(sl,'Appendix • Three high-risk viva questions','Concise answers to protect the scientific boundary',31)
card(sl,.80,1.40,3.78,4.55,'“Is this novel?”','Not as a claim for graph cascades, attribution, fault injection, or mitigation individually. The contribution is the paired integration of reliability, propagation, evidence-conditioned attribution, verification/recovery, topology, and cost.',RED,16,14)
card(sl,4.78,1.40,3.78,4.55,'“Where are the results?”','No experimental artifacts were supplied. I use a locked shell with em dashes rather than fabricate findings. This is a complete pre-experimental protocol; empirical completion follows execution and validation.',TEAL,16,14)
card(sl,8.76,1.40,3.78,4.55,'“What if H1–H6 fail?”','I report the observed direction, confidence interval, effect size, and uncertainty. Unsupported hypotheses are rejected or left inconclusive. The discussion and conclusion are rewritten accordingly.',PURPLE,16,14)
footer(sl,status=True)

out=ROOT/'MAS_RELIAB_Viva_Deck_Dev_Parth.pptx'; prs.save(out); print(out)
