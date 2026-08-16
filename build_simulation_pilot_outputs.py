from __future__ import annotations

import json
from pathlib import Path

import pandas as pd
from docx import Document
from docx.enum.section import WD_SECTION
from docx.enum.style import WD_STYLE_TYPE
from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT, WD_TABLE_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches, Pt, RGBColor
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor as PRGB
from pptx.util import Inches as PInches, Pt as PPt

ROOT = Path(__file__).resolve().parent
REPO = ROOT / "MAS-RELIAB"
TABLES = REPO / "results/tables"
ANALYSIS = REPO / "results/analysis"
FIGS = REPO / "results/figures"
METHOD_FIGS = ROOT / "figures"
PAPER = ROOT / "MAS_RELIAB_Simulation_Pilot_Paper_Dev_Parth.docx"
DECK = ROOT / "MAS_RELIAB_Simulation_Pilot_Viva_Deck_Dev_Parth.pptx"
NOTES = ROOT / "MAS_RELIAB_Simulation_Pilot_Viva_Notes_Dev_Parth.md"

NAVY = "17324D"; TEAL = "087E8B"; ORANGE = "F28C28"; LIGHT = "EAF1F5"; RED = "B23A48"; GREY = "566573"

baseline = pd.read_csv(TABLES / "table_02_baseline.csv")
attribution = pd.read_csv(TABLES / "table_05_attribution.csv")
tradeoff = pd.read_csv(TABLES / "table_08_tradeoff.csv")
tests = pd.read_csv(ANALYSIS / "statistical_tests.csv")
decisions = pd.read_csv(ANALYSIS / "hypothesis_decisions.csv")
protocol = json.loads((ANALYSIS / "analysis_manifest.json").read_text())
counts = json.loads((REPO / "results/raw/run_counts.json").read_text())

def fmt(x, digits=3):
    return f"{float(x):.{digits}f}"

def ptxt(x):
    x = float(x)
    if x == 0: return "< 10⁻³⁰⁰"
    if x < .001: return "< 0.001"
    return f"= {x:.3f}"

def test_row(h, outcome, contains):
    m = tests[(tests.hypothesis == h) & (tests.outcome == outcome) & tests.contrast.str.contains(contains, regex=False)]
    return m.iloc[0]

h1_epr = test_row("H1", "epr", "early - late")
h1_daf = test_row("H1", "daf", "early - late")
h1_depth = test_row("H1", "propagation_depth", "early - late")
h2 = test_row("H2", "top1_correct", "full_trace - output_only")
h3_epr = test_row("H3", "epr", "local;recovery=retry - verify=final")
h3_work = test_row("H3", "work_units", "local;recovery=retry - verify=final")
h4_red_rec = test_row("H4", "recovered", "redundant - recovery=retry")
h4_red_work = test_row("H4", "work_units", "redundant - recovery=retry")
h5 = test_row("H5", "daf", "middle - early")
h6 = test_row("H6", "daf", "single - parallel")
severity = test_row("RQ2-exploratory", "daf", "2 - 1")

# ---------------- DOCX ----------------
def shade(cell, fill):
    tcPr = cell._tc.get_or_add_tcPr(); shd = OxmlElement("w:shd"); shd.set(qn("w:fill"), fill); tcPr.append(shd)

def set_cell_text(cell, text, bold=False, color=None, size=8.3):
    cell.text = ""
    p = cell.paragraphs[0]; p.paragraph_format.space_after = Pt(0)
    r = p.add_run(str(text)); r.bold = bold; r.font.name = "Aptos"; r.font.size = Pt(size)
    if color: r.font.color.rgb = RGBColor.from_string(color)
    cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER

def add_table(doc, headers, rows, widths=None, font=8.2):
    table = doc.add_table(rows=1, cols=len(headers)); table.alignment = WD_TABLE_ALIGNMENT.CENTER
    table.style = "Table Grid"
    for i, h in enumerate(headers):
        set_cell_text(table.rows[0].cells[i], h, True, "FFFFFF", font); shade(table.rows[0].cells[i], NAVY)
    for ri, row in enumerate(rows):
        cells = table.add_row().cells
        for i, value in enumerate(row):
            set_cell_text(cells[i], value, False, None, font)
            if ri % 2: shade(cells[i], "F4F7F9")
            if widths: cells[i].width = Inches(widths[i])
    return table

def page_field(paragraph):
    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = paragraph.add_run(); fldChar1 = OxmlElement("w:fldChar"); fldChar1.set(qn("w:fldCharType"), "begin")
    instrText = OxmlElement("w:instrText"); instrText.set(qn("xml:space"), "preserve"); instrText.text = "PAGE"
    fldChar2 = OxmlElement("w:fldChar"); fldChar2.set(qn("w:fldCharType"), "end")
    run._r.extend([fldChar1, instrText, fldChar2])

def paragraph(doc, text="", bold_lead=None, italic=False, align=None):
    p = doc.add_paragraph(); p.paragraph_format.space_after = Pt(5); p.paragraph_format.line_spacing = 1.08
    if bold_lead and text.startswith(bold_lead):
        r = p.add_run(bold_lead); r.bold = True; p.add_run(text[len(bold_lead):])
    else: p.add_run(text)
    if italic:
        for r in p.runs: r.italic = True
    if align is not None: p.alignment = align
    return p

def bullets(doc, items):
    for item in items:
        p = doc.add_paragraph(style="List Bullet"); p.paragraph_format.space_after = Pt(3); p.add_run(item)

def figure(doc, path, caption, width=6.25):
    p = doc.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.add_run().add_picture(str(path), width=Inches(width))
    c = doc.add_paragraph(caption); c.alignment = WD_ALIGN_PARAGRAPH.CENTER
    c.style = "Caption"

def heading(doc, text, level=1):
    p = doc.add_heading(text, level=level); p.paragraph_format.keep_with_next = True; return p

def make_paper():
    d = Document()
    sec = d.sections[0]; sec.top_margin = Inches(.68); sec.bottom_margin = Inches(.68); sec.left_margin = Inches(.78); sec.right_margin = Inches(.78)
    styles = d.styles
    styles["Normal"].font.name = "Aptos"; styles["Normal"].font.size = Pt(9.5)
    styles["Title"].font.name = "Aptos Display"; styles["Title"].font.size = Pt(25); styles["Title"].font.color.rgb = RGBColor.from_string(NAVY)
    for name, size, color in [("Heading 1", 16, NAVY), ("Heading 2", 12, TEAL), ("Heading 3", 10.5, ORANGE)]:
        styles[name].font.name = "Aptos Display"; styles[name].font.size = Pt(size); styles[name].font.color.rgb = RGBColor.from_string(color)
    styles["Caption"].font.name = "Aptos"; styles["Caption"].font.size = Pt(8); styles["Caption"].font.italic = True
    for section in d.sections:
        hp = section.header.paragraphs[0]; hp.text = "MAS-RELIAB · OFFLINE SIMULATION PILOT"; hp.alignment = WD_ALIGN_PARAGRAPH.CENTER
        hp.runs[0].font.size = Pt(7.5); hp.runs[0].font.color.rgb = RGBColor.from_string(GREY)
        page_field(section.footer.paragraphs[0])

    t = d.add_paragraph(style="Title"); t.alignment = WD_ALIGN_PARAGRAPH.CENTER
    t.add_run("MAS-RELIAB")
    s = d.add_paragraph(); s.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = s.add_run("A Reproducible Simulation Pilot for Failure Propagation, Attribution, Mitigation, and Reliability–Cost Analysis in Multi-Agent Topologies")
    r.bold = True; r.font.size = Pt(15); r.font.color.rgb = RGBColor.from_string(TEAL)
    paragraph(d, "Dev Parth", align=WD_ALIGN_PARAGRAPH.CENTER).runs[0].bold = True
    paragraph(d, "Simulation-study manuscript · Version 0.2.0 · 16 August 2026", italic=True, align=WD_ALIGN_PARAGRAPH.CENTER)
    p = d.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    rr = p.add_run("EVIDENCE SCOPE: STOCHASTIC SYNTHETIC AGENTS — NOT LLM OR DEPLOYMENT EVIDENCE")
    rr.bold = True; rr.font.color.rgb = RGBColor.from_string(RED); rr.font.size = Pt(10)
    figure(d, METHOD_FIGS / "figure_01_architecture.png", "MAS-RELIAB separates benchmark methodology from an execution substrate.", 5.8)
    d.add_page_break()

    heading(d, "Abstract", 1)
    paragraph(d, f"Multi-agent systems can transform local errors into distributed failures, yet task success alone does not reveal repeatability, propagation, attribution, containment, or cost. This paper evaluates MAS-RELIAB as a unified experimental methodology in an offline stochastic simulation pilot. The implementation generates 150 objectively scored exact-state tasks in five families and measures 60 evaluation tasks across single-agent, sequential, parallel, and hierarchical topologies. Controlled wrong-value, truncation, referential-drift, and timeout faults are crossed with two severities and three functional positions. Each intervention separately records scheduling, application, consumption, lineage propagation, detection, and correction. Six experiments comprise 16,560 episodes, 82,800 trace events, and 1,800 attribution-view records. Baseline task success ranged from 0.780 to 0.870 across topologies, while pass-all across five repeats ranged from 0.333 to 0.533. Full-trace top-1 attribution reached 0.892 versus 0.302 with output-only evidence; the paired difference was 0.590 (95% CI [0.543, 0.637], Holm-adjusted p < 0.001). Local verification plus retry reduced event propagation by 0.619 relative to final-only checking plus retry, but added 0.505 simulated work units. H2 and H3 were supported; H1 received mixed support; H4–H6 were not supported. These findings establish reproducible simulator behavior and expose trade-offs, not reliability properties of real language-model agents. The main methodological limitations are oracle verification, researcher-set parameters, non-equivalent functional positions across graphs, synthetic costs, and absence of linguistic/tool environments.")
    paragraph(d, "Keywords: multi-agent systems; reliability evaluation; fault injection; failure propagation; attribution; verification; recovery; simulation")

    heading(d, "1. Introduction")
    paragraph(d, "Multi-agent architectures distribute work among specialized actors, but distribution also creates dependency paths along which a local error can be copied, transformed, amplified, or masked. A terminal task score does not identify the origin of failure, distinguish an unconsumed intervention from genuine robustness, reveal whether a mitigation interrupted a cascade, or quantify the resources used to obtain reliability. These gaps motivate an evaluation object richer than a final answer: a versioned task, an execution graph, a time-ordered trace, a documented intervention, an exact or state-based outcome, and a resource record.")
    paragraph(d, "MAS-RELIAB contributes a unified experimental methodology focused on failure propagation, attribution, mitigation, and reliability–cost trade-offs. It does not claim to be the first framework for multi-agent reliability. Related benchmarks already address coordination, repeated-run behavior, controlled failures, attribution, and cascade analysis. The contribution here is the auditable connection of these components in one executable protocol with explicit denominators and hypothesis decisions.")
    paragraph(d, "The present evidence is deliberately narrow. No hosted-model credentials or local model weights were available; therefore, the empirical study uses stochastic synthetic agents and associative exact-state tasks. This choice allows deterministic regeneration, controlled lineage tracing, and objective evaluation. It cannot reproduce natural-language ambiguity, learned behavior, prompt sensitivity, model correlation, or real tool failures. All conclusions are conditional on the simulator and configuration.")
    bullets(d, [
        "An executable repository with deterministic task generation, four topology plans, controlled fault injection, full traces, and rerun manifests.",
        "Separate reporting of intervention application, consumption, propagation, attribution, verification, recovery, and simulated resource outcomes.",
        "Six measured experiments with task-clustered confidence intervals, paired tests, effect sizes, Holm adjustment, and retained null/contradictory results.",
        "A methodological quality review that identifies what must change before any LLM or deployment claim is permissible.",
    ])

    heading(d, "2. Related Work and Positioning")
    paragraph(d, "MultiAgentBench evaluates collaboration and competition across multiple coordination protocols [2], while τ-bench emphasizes policy-constrained tool interaction and repeated-run reliability through final database state [3]. ReliabilityBench studies consistency, robustness, controlled tool/API failures, end-state equivalence, and resource use [4]. These works motivate repeated execution and state-based scoring but do not make this simulator a substitute for model-based evaluation.")
    paragraph(d, "MAST organizes observed multi-agent failure modes [1]. TraceElephant and AgenTracer focus on agent- and step-level attribution under different evidence or counterfactual conditions [5,6], while ErrorProbe combines anomaly detection, backward tracing, and executable verification [7]. AgentAsk targets information failures at handoffs [8]. MAS-RELIAB borrows the principle that attribution must be conditioned on available evidence; its full-trace method is intentionally simple and has oracle validity fields, so its accuracy is an upper-bound simulator result.")
    paragraph(d, "RiskLab supports controlled configurations for emergent-risk probing [9]; TAMAS provides adversarial safety/effectiveness tasks [10]; SILO-BENCH examines distributed coordination [11]; and From Spark to Fire models graph-based cascades and mitigation [12]. Surveys [13,14] and a unified agent-evaluation representation [15] further support separating task, environment, tool, trace, and resource records. MAS-RELIAB’s defensible novelty is therefore integration and auditability, not priority for any individual component.")
    paragraph(d, "A cell-by-cell related-work comparison is intentionally omitted from the final claims because rapidly changing preprints and conference versions make unsupported binary coverage judgments risky. Before publication, any such table should be reverified against the cited full papers rather than inferred from abstracts.")

    heading(d, "3. Research Questions and Hypotheses")
    rq_rows = [
        ("RQ1 / H6", "What does reliability reveal beyond task success?", "Rank reversals or masked weaknesses will appear across reliability dimensions."),
        ("RQ2 / H1", "How do fault type, severity, and position affect propagation?", "Earlier injections will create greater normalized impact and depth than late injections."),
        ("RQ3 / H2", "How does trace observability affect attribution?", "Full traces will improve attribution over sparse and output-only views."),
        ("RQ4 / H3", "Does local verification contain cascades?", "Local verification will reduce EPR and DAF versus final-only/no checking."),
        ("RQ5 / H4", "Which recovery strategy offers the best reliability–cost trade-off?", "Isolate and redundant recovery will exceed retry, at added cost."),
        ("RQ6 / H5", "How does topology interact with fault location?", "Parallel worker faults will be better contained; hierarchical manager faults will have larger impact than leaf faults."),
    ]
    add_table(d, ["Mapping", "Research question", "Directional hypothesis"], rq_rows, [1.0, 2.5, 3.2], 8.2)
    paragraph(d, "Hypotheses are evaluated as written. Mixed, null, and contradictory effects are retained rather than post-hoc reworded.")

    heading(d, "4. Method")
    heading(d, "4.1 Evaluation object and simulator", 2)
    paragraph(d, "An episode is represented as ⟨τ, G, Π, I, X, Y, C⟩: versioned task τ; directed acyclic execution graph G; fixed simulator policy Π; null or fault intervention I; time-ordered events X; exact outcome Y; and resource record C. Nodes combine parent artifacts and assigned task partials. Intrinsic errors occur stochastically as a function of task difficulty, topology, and node role. A unique intervention lineage remains active only while the downstream state differs from the oracle; correction or algebraic masking terminates that lineage.")
    figure(d, METHOD_FIGS / "figure_02_graph.png", "Figure 1. Episode representation and lineage-aware trace model.")
    heading(d, "4.2 Dataset and exact-state scoring", 2)
    paragraph(d, "The generator creates 150 tasks: 30 development, 30 pilot, and 90 evaluation. The measured study uses the first 60 evaluation tasks selected before execution. Five balanced families cover signed arithmetic pipelines, shard-count aggregation, evidence-vector synthesis, constraint-bitmask union, and state-delta workflows. Gold states use associative sum, vector-sum, bitwise-OR, or dictionary-sum operators. Canonical JSON equality produces the task-success label. No language-model judge or fuzzy text match is used.")
    add_table(d, ["Family", "State operator", "Objective endpoint"], [
        ("Arithmetic pipeline", "Integer sum", "Exact signed total"), ("Shard aggregation", "Integer sum", "Exact aggregate count"),
        ("Evidence synthesis", "Vector sum", "Exact support/against vector"), ("Constraint decision", "Bitwise OR", "Exact union mask"),
        ("Stateful workflow", "Dictionary sum", "Exact final state fields")], [2.0, 1.4, 3.0])
    heading(d, "4.3 Topologies", 2)
    paragraph(d, "Single-agent and sequential plans execute four dependency stages; the former reuses one agent identity while the latter assigns distinct identities at handoffs. Parallel plans use four workers and one aggregator. Hierarchical plans use four workers, two managers, and one root. This equalizes task partials but not graph size or critical path; work and simulated latency are therefore reported rather than treated as nuisance variables.")
    figure(d, METHOD_FIGS / "figure_04_topologies.png", "Figure 2. Four evaluated execution structures.")
    heading(d, "4.4 Controlled fault protocol", 2)
    paragraph(d, "Fault operators are wrong value, truncation, referential drift, and timeout. E2 crosses severity 1–2 and early, middle, and late functional positions. Severity scales numerical and truncation magnitude; timeout remains categorical. Position is mapped to each graph and is not assumed structurally equivalent. The simulator logs scheduling, application, and consumption separately. Of 11,520 E2 interventions, 11,156 were applied and 10,940 consumed (96.84% and 94.97% of scheduled, respectively). Primary propagation analyses condition on consumption.")
    heading(d, "4.5 Metrics", 2)
    bullets(d, [
        "Task success rate (TSR): exact final-state correctness; pass-all: success in all five baseline repeats; exact-output consistency: identical output across repeats.",
        "Event propagation rate (EPR): affected downstream nodes divided by nodes reachable downstream from the origin.",
        "Downstream affected final (DAF): one when the injected lineage reaches an incorrect final state; amplification and depth count affected descendants and maximum graph distance.",
        "Attribution: top-1 origin accuracy and reciprocal rank under output-only, sparse-trace, and full-trace views.",
        "Detection and recovery: trace-recorded identification and successful corrected completion after a consumed fault.",
        "Resources: deterministic work units and DAG critical-path simulated milliseconds; neither is token, currency, energy, nor observed wall-clock cost.",
    ])
    heading(d, "4.6 Experimental design and inference", 2)
    add_table(d, ["Experiment", "Executed design", "Primary purpose"], [
        ("E1", "60 tasks × 4 topologies × 5 repeats = 1,200", "Baseline multidimensional reliability"),
        ("E2", "60 × 4 × 4 faults × 2 severities × 3 positions × 2 = 11,520", "Propagation and fault factors"),
        ("E3", "600 consumed E2 cases × 3 evidence views = 1,800 records", "Attribution observability"),
        ("E4", "60 × 4 × 2 faults × 4 components = 1,920", "Verification ablation"),
        ("E5", "60 × 4 × 4 recovery modes × 2 = 1,920", "Recovery"),
        ("E6", "Operating points derived from E4–E5 and E2 contrasts", "Topology and reliability–cost synthesis")], [0.65, 3.55, 2.35])
    paragraph(d, "Means use 2,000-replicate task-clustered bootstrap 95% confidence intervals. Blocked paired contrasts use exact McNemar/binomial sign tests for binary differences and Wilcoxon signed-rank tests otherwise. Effect size is discordant-pair direction for binary outcomes or matched rank-biserial correlation for non-binary outcomes. Holm adjustment is applied across the reported contrast family. Severity comparisons are marked exploratory.")
    paragraph(d, "The paired unit is an aggregated matched cell: E2 position pairs use task × topology × fault type × severity; E2 severity pairs use task × topology × fault type × position; E3 pairs use the same source episode; E4 uses task × topology × fault type; and E5 uses task × topology. Replicates are averaged within cells. E4/E5 are blocked but not synchronized counterfactual replays because treatment contributes to the seed. E2 has no separately executed matched clean/null-intervention arm; E1 is not a seed-synchronized E2 control. Therefore, E2 propagation values describe consumed lineages rather than paired fault-minus-clean causal effects.")
    heading(d, "4.7 Reproducibility", 2)
    paragraph(d, "The master seed is 20260816. Task and episode seeds are SHA-256-derived. Raw episodes, event traces, attribution records, configurations, dataset files, software/platform metadata, source hashes, tables, figures, and hypothesis decisions are stored in the repository. The full run command is `bash scripts/run_pilot.sh`. The manifest reports Python 3.13.14 on Linux x86-64 with two logical CPUs; no Git commit was available, which is a release-process limitation.")

    heading(d, "5. Results")
    paragraph(d, f"The run completed {counts['episode_count']:,} episodes and {counts['event_count']:,} events. Results below report simulator measurements only.")
    heading(d, "5.1 E1 — Baseline multidimensional reliability (RQ1)", 2)
    rows=[]
    for topo in ["single","sequential","parallel","hierarchical"]:
        r=baseline[baseline.topology==topo].iloc[0]
        rows.append((topo.title(), f"{fmt(r.task_success)} [{fmt(r.task_success_ci_low)}, {fmt(r.task_success_ci_high)}]", fmt(r.pass_k), fmt(r.exact_output_consistency), fmt(r.work_units,2), fmt(r.latency_ms,2)))
    add_table(d, ["Topology", "TSR [95% CI]", "Pass-all", "Consistency", "Work", "Sim. ms"], rows, font=7.8)
    paragraph(d, "Single-agent execution had the highest mean TSR (0.870), but sequential execution had the highest pass-all value (0.533 versus 0.483). Parallel execution had the shortest simulated critical-path latency (25.75 ms) but lower TSR and pass-all. Hierarchical execution used the most work units and had the lowest baseline TSR. The direct single-versus-parallel DAF contrast underlying H6 was not significant after Holm correction, so the descriptive rank difference was not elevated to a supported masked-weakness claim.")
    figure(d, FIGS / "figure_01_baseline_reliability.png", "Figure 3. Baseline reliability dimensions; bars are measured simulator means.")

    heading(d, "5.2 E2 — Fault propagation (RQ2, H1)", 2)
    add_table(d, ["Paired contrast", "Mean difference [95% CI]", "Holm p", "Effect"], [
        ("Early − late EPR", f"{fmt(h1_epr.mean_difference)} [{fmt(h1_epr.ci_low)}, {fmt(h1_epr.ci_high)}]", ptxt(h1_epr.p_adjusted_holm), fmt(h1_epr.effect_size)),
        ("Early − late depth", f"{fmt(h1_depth.mean_difference)} [{fmt(h1_depth.ci_low)}, {fmt(h1_depth.ci_high)}]", ptxt(h1_depth.p_adjusted_holm), fmt(h1_depth.effect_size)),
        ("Early − late DAF", f"{fmt(h1_daf.mean_difference)} [{fmt(h1_daf.ci_low)}, {fmt(h1_daf.ci_high)}]", ptxt(h1_daf.p_adjusted_holm), fmt(h1_daf.effect_size)),
        ("Severity 2 − 1 DAF", f"{fmt(severity.mean_difference)} [{fmt(severity.ci_low)}, {fmt(severity.ci_high)}]", ptxt(severity.p_adjusted_holm), fmt(severity.effect_size)),
    ], font=7.9)
    paragraph(d, "Early injection increased EPR by 0.969 and depth by 2.181 because more descendants were reachable. In contrast, early DAF was 0.020 lower than late DAF: a late corruption can directly fail the final state even though it has no downstream descendants. H1 therefore received mixed support. Severity level did not significantly change pooled DAF. This is a structural warning against interpreting EPR without DAF and reachability.")
    figure(d, FIGS / "figure_02_fault_propagation.png", "Figure 4. Propagation by functional position and topology; late EPR is structurally zero because no downstream nodes remain.")

    heading(d, "5.3 E3 — Observability-conditioned attribution (RQ3, H2)", 2)
    rows=[]
    for view in ["output_only","sparse_trace","full_trace"]:
        r=attribution[attribution.observability==view].iloc[0]
        rows.append((view.replace("_"," ").title(), f"{fmt(r.top1_correct)} [{fmt(r.top1_correct_ci_low)}, {fmt(r.top1_correct_ci_high)}]", f"{fmt(r.reciprocal_rank)} [{fmt(r.reciprocal_rank_ci_low)}, {fmt(r.reciprocal_rank_ci_high)}]"))
    add_table(d, ["Evidence view", "Top-1 accuracy [95% CI]", "MRR [95% CI]"], rows)
    paragraph(d, f"Full traces improved top-1 attribution over output-only evidence by {fmt(h2.mean_difference)} (95% CI [{fmt(h2.ci_low)}, {fmt(h2.ci_high)}], Holm p {ptxt(h2.p_adjusted_holm)}, effect {fmt(h2.effect_size)}). H2 was supported. However, full traces expose oracle local-validity information; this quantifies information availability in the simulator rather than the accuracy of a production debugger.")
    figure(d, FIGS / "figure_03_attribution.png", "Figure 5. Attribution accuracy under three trace-evidence views.")

    heading(d, "5.4 E4 — Verification ablation (RQ4, H3)", 2)
    paragraph(d, f"Local verification plus retry reduced EPR by {abs(h3_epr.mean_difference):.3f} relative to final-only verification plus retry (local − final {fmt(h3_epr.mean_difference)}, 95% CI [{fmt(h3_epr.ci_low)}, {fmt(h3_epr.ci_high)}], Holm p {ptxt(h3_epr.p_adjusted_holm)}, effect {fmt(h3_epr.effect_size)}). It added {fmt(h3_work.mean_difference)} work units (95% CI [{fmt(h3_work.ci_low)}, {fmt(h3_work.ci_high)}]). H3 was supported: early checking interrupted lineage before aggregation, but did so using privileged exact-state verification.")
    figure(d, FIGS / "figure_04_verification.png", "Figure 6. Verification ablation: reliability improvement and containment vary with topology and checking location.")

    heading(d, "5.5 E5 — Recovery strategies (RQ5, H4)", 2)
    paragraph(d, f"Redundant recovery had a positive raw recovery-rate difference versus retry ({fmt(h4_red_rec.mean_difference)}, 95% CI [{fmt(h4_red_rec.ci_low)}, {fmt(h4_red_rec.ci_high)}]) but was not significant after Holm correction (p {ptxt(h4_red_rec.p_adjusted_holm)}). Its work cost was {fmt(h4_red_work.mean_difference)} units higher (95% CI [{fmt(h4_red_work.ci_low)}, {fmt(h4_red_work.ci_high)}], Holm p {ptxt(h4_red_work.p_adjusted_holm)}). Isolate recovery was cheaper but significantly less effective than retry. The joint claim that isolate and redundant recovery outperform retry was not supported; H4 was rejected as formulated.")
    figure(d, FIGS / "figure_05_recovery.png", "Figure 7. Recovery rate and simulated work-unit trade-off.")

    heading(d, "5.6 E6 — Topology and reliability–cost synthesis (RQ6, H5–H6)", 2)
    paragraph(d, f"For H5, hierarchical manager-position DAF exceeded leaf-position DAF by only {fmt(h5.mean_difference)} (95% CI [{fmt(h5.ci_low)}, {fmt(h5.ci_high)}], Holm p {ptxt(h5.p_adjusted_holm)}), and the proposed parallel-worker containment contrast was also non-significant. H5 was not supported. For H6, task-success and DAF-containment ranks differed descriptively, but single-minus-parallel DAF was {fmt(h6.mean_difference)} (95% CI [{fmt(h6.ci_low)}, {fmt(h6.ci_high)}], Holm p {ptxt(h6.p_adjusted_holm)}). H6 was not supported by its direct inferential contrast.")
    paragraph(d, "The operating points nevertheless demonstrate why components should remain separate: parallel plans can shorten critical-path latency while increasing work, and redundant recovery can increase success while moving to a higher-cost region. No scalar composite is reported because arbitrary weights would encode an undeclared utility function.")
    figure(d, FIGS / "figure_06_tradeoff.png", "Figure 8. Reliability–cost operating points. Work and latency are simulator quantities, not deployment measurements.")

    heading(d, "5.7 Hypothesis decision summary", 2)
    add_table(d, ["Hypothesis", "Decision", "Observed basis"], [(r.hypothesis, r.decision, r.observed_basis) for _,r in decisions.iterrows()], [0.65, 1.15, 4.7], 7.7)

    heading(d, "6. Discussion")
    heading(d, "6.1 Answers to RQ1–RQ6", 2)
    paragraph(d, "RQ1: Baseline TSR alone hid differences in pass-all, consistency, work, and simulated latency, but the prespecified H6 inferential criterion was not met. The correct conclusion is descriptive multidimensionality, not a statistically established rank reversal.")
    paragraph(d, "RQ2: Position strongly changed cascade opportunity. Early faults affected more downstream nodes and produced greater depth; late faults directly damaged the terminal state and therefore had slightly higher DAF. Severity did not have a pooled DAF effect. Reporting EPR without DAF would have reversed part of the interpretation.")
    paragraph(d, "RQ3: Trace evidence substantially improved origin localization under a fixed heuristic. The result supports observability as an experimental factor, while the oracle validity field limits realism.")
    paragraph(d, "RQ4: Local oracle checking was the strongest containment component and increased work. The result is an upper bound on the value of accurate intermediate verification.")
    paragraph(d, "RQ5: Retry was difficult to dominate. Redundancy cost more and its recovery gain did not survive multiplicity correction; isolation saved work but recovered fewer faults. No universal winner emerged.")
    paragraph(d, "RQ6: Proposed topology-by-location patterns were not detected. Graph size and reachable-set geometry explained visible EPR/depth differences, but normalized final impact was similar. Topology recommendations would be premature.")
    heading(d, "6.2 Implications for MAS-RELIAB", 2)
    paragraph(d, "The pilot validates the mechanics of an auditable methodology: intervention integrity can be separated from propagation; trace views can be held fixed on the same failed episode; mitigations can be evaluated jointly with cost; and hypotheses can fail without invalidating the measurement pipeline. It also demonstrates why a composite reliability score should not precede component metrics. Stakeholders may value success, latency, cost, containment, and attribution differently; those preferences should be explicit utility functions, not hidden scientific weights.")
    paragraph(d, "The strongest next use of the repository is preregistration support. A model-based replication can preserve task IDs, graph schemas, fault IDs, pair blocks, metric code, and decision rules while replacing synthetic actors, oracle verifiers, and simulated resources with real prompts, tools, traces, and costs.")

    heading(d, "7. Threats to Validity and Methodological Review")
    heading(d, "7.1 Construct validity", 2)
    paragraph(d, "Associative state tasks are objectively checkable but do not represent open-ended reasoning. EPR and DAF capture explicit lineage effects, not semantic influence. Exact oracle verification is stronger than realistic validation. Functional positions differ structurally by topology, and timeout severity is categorical rather than ordinal.")
    heading(d, "7.2 Internal validity", 2)
    paragraph(d, "Researcher-set probabilities govern intrinsic errors, detection, correction, and protocol application. The study uses task blocking but not fully separated common-random-number streams for every treatment component. A treatment can consume additional random draws and alter later stochastic events. Future code should allocate independent RNG streams to intrinsic errors, intervention, detection, recovery, and latency and run null counterfactuals with synchronized streams.")
    heading(d, "7.3 Statistical conclusion validity", 2)
    paragraph(d, "The pilot has many episodes but only 60 measured tasks from five generated families. Task-clustered intervals reduce pseudoreplication, and Holm correction addresses the reported contrast family, but family-level generalization remains weak. H5 and H6 may also be underpowered with four topologies and small normalized-impact differences. Replication should increase independent task templates and use hierarchical task-family models.")
    heading(d, "7.4 External validity", 2)
    paragraph(d, "No LLM, natural-language prompt, tool API, model provider, production scheduler, or user interaction was evaluated. Tokens are not applicable; latency and work are simulated. The study cannot support claims about actual agents, frameworks, vendors, safety, or deployment reliability. A full prioritized review and remedies are provided in `docs/METHODOLOGY_REVIEW.md`.")

    heading(d, "8. Conclusion")
    paragraph(d, "MAS-RELIAB was implemented and executed as a reproducible offline simulation pilot. Across 16,560 episodes, full traces improved attribution and local oracle verification contained cascades, while earlier faults created deeper, broader propagation but not greater final impact. Several anticipated claims failed: recovery alternatives did not jointly dominate retry, topology-by-location effects were not supported, and the direct rank-discordance contrast did not support H6. These mixed findings are the result, not a defect to hide.")
    paragraph(d, "The empirical contribution is limited to a configured synthetic environment. The repository demonstrates controlled intervention, trace lineage, state-based evaluation, paired analysis, effect reporting, and explicit reliability–cost accounting. Establishing relevance to language-model multi-agent systems requires an external replication with real models, tools, fallible verifiers, observed costs, broader tasks, synchronized counterfactual streams, and preregistered sensitivity analysis.")

    heading(d, "Data and Code Availability")
    paragraph(d, "The executable artifact is in the accompanying `MAS-RELIAB/` repository. Generated tasks, raw episodes, event traces, configurations, seeds, manifests, measured tables, figures, and tests are included. Citation metadata and the MIT License are provided. The placeholder repository URL in `CITATION.cff` must be replaced before public release.")

    heading(d, "References")
    refs = [
        '[1] M. Cemri et al., “Why Do Multi-Agent LLM Systems Fail?” arXiv:2503.13657, 2025.',
        '[2] K. Zhu et al., “MultiAgentBench: Evaluating the Collaboration and Competition of LLM Agents,” arXiv:2503.01935, 2025.',
        '[3] S. Yao, N. Shinn, P. Razavi, and K. Narasimhan, “τ-bench: A Benchmark for Tool-Agent-User Interaction in Real-World Domains,” arXiv:2406.12045, 2024.',
        '[4] A. Gupta, “ReliabilityBench: Evaluating LLM Agent Reliability Under Production-Like Stress Conditions,” arXiv:2601.06112, 2026.',
        '[5] M. Chen et al., “Seeing the Whole Elephant: A Benchmark for Failure Attribution in LLM-based Multi-Agent Systems,” ACL 2026; arXiv:2604.22708.',
        '[6] G. Zhang et al., “AgenTracer: Who Is Inducing Failure in the LLM Agentic Systems?” arXiv:2509.03312, 2025.',
        '[7] J. Li, E. Yilmaz, B. Chen, and T. Le, “Towards Self-Improving Error Diagnosis in Multi-Agent Systems,” Findings of ACL 2026, pp. 2063–2077.',
        '[8] B. Lin et al., “AgentAsk: Multi-Agent Systems Need to Ask,” ACL 2026, pp. 28055–28077.',
        '[9] Y. Jiang et al., “RiskLab: A Controlled Toolkit for Probing Emergent Risks in LLM-Based Multi-Agent Systems,” ACL 2026 System Demonstrations, pp. 167–177.',
        '[10] I. Kavathekar et al., “TAMAS: Benchmarking Adversarial Risks in Multi-Agent LLM Systems,” arXiv:2511.05269, 2025.',
        '[11] Y. Zhang et al., “SILO-BENCH: A Scalable Environment for Evaluating Distributed Coordination in Multi-Agent LLM Systems,” ACL 2026; arXiv:2603.01045.',
        '[12] Y. Xie et al., “From Spark to Fire: Modeling and Mitigating Error Cascades in LLM-Based Multi-Agent Collaboration,” arXiv:2603.04474, 2026.',
        '[13] A. Yehudai et al., “Survey on Evaluation of LLM-based Agents,” arXiv:2503.16416, 2025.',
        '[14] M. Mohammadi et al., “Evaluation and Benchmarking of LLM Agents: A Survey,” Proc. ACM SIGKDD, 2025; arXiv:2507.21504.',
        '[15] P. Zhu et al., “A Unified Framework for the Evaluation of LLM Agentic Capabilities,” arXiv:2605.27898, 2026.',
        '[16] B. Efron and R. Tibshirani, An Introduction to the Bootstrap. Chapman & Hall/CRC, 1993.',
        '[17] S. Holm, “A Simple Sequentially Rejective Multiple Test Procedure,” Scandinavian Journal of Statistics, 1979.',
    ]
    for ref in refs:
        p=paragraph(d, ref); p.paragraph_format.left_indent=Inches(.2); p.paragraph_format.first_line_indent=Inches(-.2); p.runs[0].font.size=Pt(8)

    heading(d, "Appendix A. Reproduction Checklist")
    bullets(d, [
        "Install package dependencies from `pyproject.toml` or `requirements.txt`.",
        "Run `bash scripts/run_pilot.sh`; it regenerates data, runs experiments, analyzes results, executes tests, and verifies required artifacts.",
        "Confirm `results/raw/reproducibility_manifest.json` hashes and environment fields.",
        "Use `results/analysis/statistical_tests.csv` rather than values copied from this manuscript for machine analysis.",
        "Do not describe these outputs as LLM-agent or deployment measurements.",
    ])
    heading(d, "Appendix B. Hypothesis Integrity Matrix")
    add_table(d, ["Hypothesis", "Primary contrast", "Decision", "Integrity note"], [
        ("H1", "Early vs late EPR/depth/DAF", "Mixed", "DAF contradicted EPR/depth direction"),
        ("H2", "Full trace vs output only", "Supported", "Oracle trace limits realism"),
        ("H3", "Local vs final/no verification", "Supported", "Exact oracle upper bound"),
        ("H4", "Isolate/redundant vs retry", "Not supported", "Raw redundant gain failed Holm correction"),
        ("H5", "Topology × functional location", "Not supported", "No significant planned pattern"),
        ("H6", "TSR leader vs DAF-containment leader", "Not supported", "Descriptive ranking alone was insufficient"),
    ], font=7.7)
    d.save(PAPER)

# ---------------- PPTX ----------------
PNAVY=PRGB(23,50,77); PTEAL=PRGB(8,126,139); PORANGE=PRGB(242,140,40); PWHITE=PRGB(255,255,255); PLIGHT=PRGB(234,241,245); PRED=PRGB(178,58,72); PGREY=PRGB(86,101,115)

def rect(slide,x,y,w,h,color, radius=False):
    shape=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,PInches(x),PInches(y),PInches(w),PInches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb=color; shape.line.fill.background(); return shape

def textbox(slide,text,x,y,w,h,size=20,color=PNAVY,bold=False,align=PP_ALIGN.LEFT,margin=.08):
    box=slide.shapes.add_textbox(PInches(x),PInches(y),PInches(w),PInches(h)); tf=box.text_frame; tf.clear(); tf.word_wrap=True
    tf.margin_left=tf.margin_right=PInches(margin); tf.margin_top=tf.margin_bottom=PInches(margin)
    p=tf.paragraphs[0]; p.alignment=align; r=p.add_run(); r.text=text; r.font.name="Aptos"; r.font.size=PPt(size); r.font.bold=bold; r.font.color.rgb=color
    return box

def bulletbox(slide,items,x,y,w,h,size=18,color=PNAVY):
    box=slide.shapes.add_textbox(PInches(x),PInches(y),PInches(w),PInches(h)); tf=box.text_frame; tf.clear(); tf.word_wrap=True
    tf.margin_left=PInches(.1); tf.margin_right=PInches(.06); tf.margin_top=PInches(.05)
    for i,item in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.text=item; p.level=0; p.font.name="Aptos"; p.font.size=PPt(size); p.font.color.rgb=color; p.space_after=PPt(9); p.text="• "+p.text
    return box

def base_slide(prs,title,kicker=None):
    s=prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.333,.18,PTEAL); textbox(s,title,.55,.34,12.1,.5,26,PNAVY,True)
    if kicker: textbox(s,kicker,.58,.86,11.9,.3,10,PGREY,False)
    rect(s,.55,7.13,12.2,.02,PLIGHT); textbox(s,"MAS-RELIAB · Dev Parth · Offline simulation pilot",.58,7.18,7,.18,8,PGREY)
    textbox(s,str(len(prs.slides)),12.2,7.16,.45,.2,8,PGREY,False,PP_ALIGN.RIGHT)
    return s

def picture(slide,path,x,y,w,h=None):
    if h: return slide.shapes.add_picture(str(path),PInches(x),PInches(y),PInches(w),PInches(h))
    return slide.shapes.add_picture(str(path),PInches(x),PInches(y),width=PInches(w))

def callout(slide,big,label,x,y,w,color=PTEAL):
    rect(slide,x,y,w,1.05,PLIGHT,True); textbox(slide,big,x+.08,y+.09,w-.16,.42,25,color,True,PP_ALIGN.CENTER); textbox(slide,label,x+.08,y+.58,w-.16,.28,10,PGREY,False,PP_ALIGN.CENTER)

def make_deck():
    prs=Presentation(); prs.slide_width=PInches(13.333); prs.slide_height=PInches(7.5)
    s=prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.333,7.5,PNAVY); rect(s,0,0,.23,7.5,PTEAL)
    textbox(s,"MAS-RELIAB",.75,1.0,11.6,.65,34,PWHITE,True); textbox(s,"Measured offline simulation pilot",.78,1.72,11.5,.45,22,PTEAL,True)
    textbox(s,"Failure propagation · attribution · mitigation · reliability–cost trade-offs",.78,2.36,11.2,.65,22,PWHITE)
    textbox(s,"Dev Parth",.78,4.45,5,.4,18,PWHITE,True); textbox(s,"16 August 2026 · v0.2.0",.78,4.92,5,.3,12,PRGB(190,205,215))
    rect(s,.78,5.75,10.9,.62,PRED,True); textbox(s,"Synthetic-agent evidence — not LLM or deployment evidence",.95,5.91,10.55,.27,17,PWHITE,True,PP_ALIGN.CENTER)

    s=base_slide(prs,"Evidence boundary","The most important interpretation rule")
    callout(s,"0","LLM/API calls",.7,1.45,2.5,PRED); callout(s,"150","generated tasks",3.45,1.45,2.5); callout(s,"16,560","episodes",6.2,1.45,2.5); callout(s,"82,800","trace events",8.95,1.45,2.5)
    bulletbox(s,["Validates implementation, causal bookkeeping, and configured simulator mechanisms.","Does not establish real-agent reliability, safety, or topology superiority.","All cost and latency values are simulated; tokens are not applicable."],.8,3.05,11.5,2.7,19)

    s=base_slide(prs,"Problem: success alone is diagnostically weak")
    picture(s,METHOD_FIGS/"figure_02_graph.png",.65,1.3,6.2)
    bulletbox(s,["Was an injected fault applied and consumed?","Where did it originate and how far did it spread?","Could available evidence attribute it?","Did verification or recovery stop it?","What reliability cost was paid?"],7.05,1.45,5.45,4.8,18)

    s=base_slide(prs,"Contribution and positioning")
    rect(s,.7,1.25,5.7,4.9,PLIGHT,True); textbox(s,"What MAS-RELIAB contributes",.95,1.52,5.1,.35,21,PTEAL,True)
    bulletbox(s,["One auditable methodology linking controlled faults, lineage, attribution, mitigation, and cost.","Exact-state tasks and explicit intervention denominators.","Executable repository, raw traces, paired inference, and negative findings."],.95,2.05,5.05,3.5,17)
    rect(s,6.8,1.25,5.7,4.9,PLIGHT,True); textbox(s,"What it does not claim",7.05,1.52,5.1,.35,21,PRED,True)
    bulletbox(s,["Not the first multi-agent reliability framework.","Not a replacement for MAST, MultiAgentBench, ReliabilityBench, TraceElephant, RiskLab, or cascade benchmarks.","Not evidence about any model, vendor, or production system."],7.05,2.05,5.05,3.5,17)

    s=base_slide(prs,"Research questions and hypotheses")
    rows=[("RQ1/H6","Multidimensional reliability"),("RQ2/H1","Fault type, severity, position"),("RQ3/H2","Observability and attribution"),("RQ4/H3","Local verification"),("RQ5/H4","Recovery–cost trade-off"),("RQ6/H5","Topology × location")]
    for i,(a,b) in enumerate(rows):
        x=.75+(i%2)*6.15; y=1.25+(i//2)*1.72; rect(s,x,y,5.72,1.28,PLIGHT,True); textbox(s,a,x+.18,y+.16,1.1,.3,15,PTEAL,True); textbox(s,b,x+1.25,y+.15,4.2,.62,18,PNAVY,True)
    textbox(s,"Hypotheses were decisions, not promises: mixed and null results were retained.",1.1,6.45,11.1,.35,16,PRED,True,PP_ALIGN.CENTER)

    s=base_slide(prs,"Simulator and exact-state tasks")
    picture(s,METHOD_FIGS/"figure_01_architecture.png",.7,1.25,6.2)
    bulletbox(s,["Five task families; associative exact-state gold labels.","Four topology plans with stochastic intrinsic errors.","Unique fault lineage follows corrupted state through the DAG.","Lineage stops on correction or algebraic masking.","No subjective judge or text similarity."],7.05,1.42,5.3,4.8,17)

    s=base_slide(prs,"Four execution topologies")
    picture(s,METHOD_FIGS/"figure_04_topologies.png",.75,1.22,7.1)
    bulletbox(s,["Single: one identity across four stages","Sequential: four-agent handoff chain","Parallel: four workers + aggregator","Hierarchical: workers + managers + root","Graph size and critical path are not equal — cost is reported"],8.05,1.42,4.6,4.7,16)

    s=base_slide(prs,"Controlled fault protocol")
    for i,(name,desc) in enumerate([("Wrong value","state perturbation"),("Truncation","component loss"),("Referential drift","bit/field displacement"),("Timeout","missing state")]):
        x=.75+i*3.08; rect(s,x,1.35,2.72,1.3,PLIGHT,True); textbox(s,name,x+.1,1.56,2.52,.3,18,PTEAL,True,PP_ALIGN.CENTER); textbox(s,desc,x+.1,2.05,2.52,.25,11,PGREY,False,PP_ALIGN.CENTER)
    bulletbox(s,["Severity: 1–2 · Position: early / middle / late","Scheduled ≠ applied ≠ consumed — each is logged","11,520 E2 faults; 96.84% applied; 94.97% consumed","Propagation analyses condition on consumed faults"],1.15,3.15,11.0,2.65,18)

    s=base_slide(prs,"Six experiments and component analysis")
    items=[("E1","Baseline reliability","1,200 eps"),("E2","Fault propagation","11,520 eps"),("E3","Attribution views","1,800 records"),("E4","Verification ablation","1,920 eps"),("E5","Recovery","1,920 eps"),("E6","Cost/topology synthesis","derived")]
    for i,(eid,name,n) in enumerate(items):
        x=.72+(i%3)*4.2; y=1.25+(i//3)*2.35; rect(s,x,y,3.82,1.83,PLIGHT,True); textbox(s,eid,x+.18,y+.17,.65,.35,22,PORANGE,True); textbox(s,name,x+.88,y+.2,2.7,.55,18,PNAVY,True); textbox(s,n,x+.88,y+.95,2.65,.3,14,PTEAL,True)
    textbox(s,"Inference: task-clustered bootstrap CIs · paired exact/Wilcoxon tests · effect sizes · Holm adjustment",.85,6.15,11.7,.45,15,PGREY,True,PP_ALIGN.CENTER)

    s=base_slide(prs,"E1 result: reliability is multidimensional")
    picture(s,FIGS/"figure_01_baseline_reliability.png",.55,1.18,8.15)
    callout(s,"0.870","highest TSR: single",9.05,1.45,3.2); callout(s,"0.533","highest pass-all: sequential",9.05,2.85,3.2,PORANGE); callout(s,"25.75 ms","shortest sim. latency: parallel",9.05,4.25,3.2,PTEAL)
    textbox(s,"H6 not supported by the direct paired DAF contrast.",8.9,5.75,3.5,.55,16,PRED,True,PP_ALIGN.CENTER)

    s=base_slide(prs,"E2 result: position changes opportunity and final impact differently")
    picture(s,FIGS/"figure_02_fault_propagation.png",.55,1.18,8.2)
    callout(s,"+0.969","early − late EPR",9.05,1.45,3.25); callout(s,"+2.181","early − late depth",9.05,2.85,3.25); callout(s,"−0.020","early − late DAF",9.05,4.25,3.25,PRED)
    textbox(s,"H1: MIXED SUPPORT",9.0,5.72,3.3,.45,18,PRED,True,PP_ALIGN.CENTER)

    s=base_slide(prs,"E3 result: trace evidence improves attribution")
    picture(s,FIGS/"figure_03_attribution.png",.65,1.23,7.05)
    callout(s,"0.302","output only",8.15,1.35,2.0,PRED); callout(s,"0.803","sparse trace",10.35,1.35,2.0,PORANGE); callout(s,"0.892","full trace",9.25,2.78,2.0,PTEAL)
    textbox(s,"Paired full − output: +0.590\n95% CI [0.543, 0.637]\nHolm p < 0.001 · effect 0.839",8.15,4.2,4.2,1.2,17,PNAVY,True,PP_ALIGN.CENTER)
    textbox(s,"H2 SUPPORTED — but full trace contains oracle validity",8.05,5.78,4.5,.5,15,PRED,True,PP_ALIGN.CENTER)

    s=base_slide(prs,"E4 result: local oracle verification contains cascades")
    picture(s,FIGS/"figure_04_verification.png",.48,1.18,8.35)
    callout(s,"−0.619","local − final EPR",9.15,1.52,3.1,PTEAL); callout(s,"+0.505","work units",9.15,3.05,3.1,PORANGE)
    textbox(s,"Both Holm p < 0.001",9.2,4.55,3.0,.35,15,PNAVY,True,PP_ALIGN.CENTER)
    textbox(s,"H3 SUPPORTED\nUpper-bound result under privileged exact-state checking",8.85,5.2,3.7,.95,16,PRED,True,PP_ALIGN.CENTER)

    s=base_slide(prs,"E5 result: no universal recovery winner")
    picture(s,FIGS/"figure_05_recovery.png",.48,1.18,8.35)
    callout(s,"+0.036","redundant − retry recovery",9.05,1.35,3.3,PTEAL); callout(s,"+1.013","redundant − retry work",9.05,2.8,3.3,PORANGE)
    textbox(s,"Recovery gain failed Holm correction (p = 0.356).\nIsolation was cheaper but recovered fewer faults.",8.95,4.25,3.5,1.0,16,PNAVY,True,PP_ALIGN.CENTER)
    textbox(s,"H4 NOT SUPPORTED",9.1,5.72,3.2,.45,18,PRED,True,PP_ALIGN.CENTER)

    s=base_slide(prs,"E6 result: reliability–cost operating points")
    picture(s,FIGS/"figure_06_tradeoff.png",.35,1.18,8.7)
    bulletbox(s,["Redundancy moves upward and right: more success, more work.","Parallel plans shorten simulated critical path but use more work.","No scalar composite — weights would encode stakeholder utility.","H5 and H6 were not supported."],9.05,1.55,3.65,3.8,15)

    s=base_slide(prs,"Hypothesis decisions: report what happened")
    color_map={"supported":PTEAL,"mixed support":PORANGE,"not supported":PRED}
    for i,row in decisions.iterrows():
        y=1.2+i*.83; textbox(s,row.hypothesis,.8,y,.65,.35,19,PNAVY,True); rect(s,1.55,y-.03,2.0,.55,color_map[row.decision],True); textbox(s,row.decision.upper(),1.62,y+.08,1.85,.24,12,PWHITE,True,PP_ALIGN.CENTER); textbox(s,row.observed_basis,3.8,y-.02,8.65,.55,13,PNAVY)
    textbox(s,"Negative and contradictory outcomes were not converted into confirmations.",1.2,6.45,10.9,.35,16,PRED,True,PP_ALIGN.CENTER)

    s=base_slide(prs,"Methodological weaknesses — research-quality review")
    bulletbox(s,["No external validity: synthetic state transformers, not LLM agents.","Oracle intermediate states make verification and attribution optimistic.","Researcher-set rates can determine rankings — sensitivity analysis is required.","Functional positions are not graph-equivalent; late EPR is structurally zero.","60 measured tasks and five algebraic families limit generalization.","Work and latency are simulated, not observed cost."],.85,1.25,7.25,4.9,17)
    rect(s,8.45,1.35,3.9,4.7,PLIGHT,True); textbox(s,"Verdict",8.8,1.7,3.2,.4,25,PTEAL,True,PP_ALIGN.CENTER); textbox(s,"Fit for",8.75,2.5,3.3,.3,17,PNAVY,True,PP_ALIGN.CENTER); textbox(s,"implementation pilot\nmethod validation\nhypothesis refinement",8.8,2.9,3.2,1.15,16,PNAVY,False,PP_ALIGN.CENTER); textbox(s,"Not fit for",8.75,4.32,3.3,.3,17,PRED,True,PP_ALIGN.CENTER); textbox(s,"deployment claims\nmodel rankings\nsafety certification",8.8,4.72,3.2,1.05,16,PRED,False,PP_ALIGN.CENTER)

    s=base_slide(prs,"Reproducible artifact")
    picture(s,METHOD_FIGS/"figure_05_experiments.png",.55,1.25,6.4)
    bulletbox(s,["`bash scripts/run_pilot.sh` regenerates data and results.","Raw episodes, 82,800 events, seeds, config, source/data hashes.","Eight measured CSV tables + six empirical figures.","18 tests pass; GitHub Actions workflow included.","Data card, fault protocol, experiment registry, methodology review.","Release weakness: manifest has no Git commit; tag before publication."],7.1,1.35,5.35,4.75,16)

    s=base_slide(prs,"What the pilot establishes — and what comes next")
    rect(s,.75,1.25,5.8,4.85,PLIGHT,True); textbox(s,"Established inside simulator",1.0,1.55,5.3,.35,21,PTEAL,True)
    bulletbox(s,["Auditable intervention and lineage mechanics","Trace evidence changes attribution","Local oracle verification contains cascades","Recovery has measurable reliability–cost trade-offs","Hypotheses can be rejected transparently"],1.0,2.05,5.15,3.5,16)
    rect(s,6.8,1.25,5.8,4.85,PLIGHT,True); textbox(s,"Required external replication",7.05,1.55,5.3,.35,21,PORANGE,True)
    bulletbox(s,["Multiple real model families and frozen prompts","Sandboxed tools and broader task domains","Fallible verifiers and realistic logs","Independent synchronized RNG/counterfactual streams","Observed tokens, latency, cost, hardware","Preregistered parameter sensitivity"],7.05,2.05,5.15,3.7,16)

    s=base_slide(prs,"Conclusion")
    textbox(s,"MAS-RELIAB is now executable, measured, and auditable.",.85,1.25,11.6,.55,25,PTEAL,True,PP_ALIGN.CENTER)
    bulletbox(s,["16,560 episodes across four topologies and six experiments.","H2 and H3 supported; H1 mixed; H4–H6 not supported.","Report EPR with DAF and reachability — no single score is sufficient.","The honest contribution is the integrated experimental methodology.","The honest limitation is equally clear: this is simulation-pilot evidence only."],1.25,2.15,10.8,3.45,20)
    textbox(s,"Questions?",4.95,6.05,3.4,.55,28,PNAVY,True,PP_ALIGN.CENTER)

    s=base_slide(prs,"Viva backup: likely questions")
    bulletbox(s,["Why simulate instead of use LLMs?","Why is H1 mixed rather than supported?","Why is late EPR zero but DAF high?","Does full-trace attribution leak the answer?","Why not publish a composite reliability score?","How would you establish external validity?","What would falsify the methodology’s usefulness?"],1.0,1.25,11.3,4.9,18)
    textbox(s,"Answer pattern: result → limitation → next validation step",1.8,6.25,9.7,.4,18,PTEAL,True,PP_ALIGN.CENTER)

    prs.save(DECK)


def make_notes():
    slides = [
        ("Title", "Open with the evidence boundary: this is the measured simulation pilot, not an LLM benchmark."),
        ("Evidence boundary", "State the scale and what zero model calls means. Treat this as integrity, not an apology."),
        ("Problem", "Use one example: final failure alone cannot reveal whether the injection failed, propagated, or was detected."),
        ("Contribution", "Avoid priority language. Emphasize integration and auditability."),
        ("RQs", "Explain that hypotheses were frozen directional claims; null findings remain informative."),
        ("Simulator", "Explain exact-state algebra and lineage termination. Mention stochastic intrinsic errors."),
        ("Topologies", "Clarify that graph size differs and therefore cost/critical path must be shown."),
        ("Fault protocol", "Explain scheduled/applied/consumed denominators. Propagation uses consumed faults."),
        ("Experiments", "E3 and E6 are derived analyses, so episode totals do not double-count them."),
        ("E1", "Point out that single leads TSR but sequential leads pass-all. H6 still failed its direct inferential criterion."),
        ("E2", "This is the key counterintuitive result: early creates more cascade opportunity; late directly breaks final state."),
        ("E3", "State the +0.590 paired effect, then immediately disclose oracle-validity optimism."),
        ("E4", "State containment gain and cost together. Call it an upper bound."),
        ("E5", "Do not claim redundancy wins; its raw recovery gain failed Holm correction."),
        ("E6", "Explain the Pareto idea without claiming a universal optimum or composite score."),
        ("Decisions", "Read H1, H4, H5, H6 honestly. Stress that the protocol survives rejected hypotheses."),
        ("Weaknesses", "Prioritize external validity, oracle verification, parameter sensitivity, and position inequivalence."),
        ("Artifact", "Mention raw trace count, hashes, tests, CI, and missing Git commit as an openly recorded weakness."),
        ("Next study", "Propose real models, tools, fallible verifiers, synchronized counterfactual streams, observed costs."),
        ("Conclusion", "Finish with the integrated methodology and evidence boundary."),
        ("Backup", "Use result → limitation → next step for every challenge."),
    ]
    lines=["# MAS-RELIAB simulation-pilot viva notes", "", "Target talk: 14–16 minutes plus questions.", ""]
    for i,(title,note) in enumerate(slides,1): lines += [f"## Slide {i}: {title}", note, ""]
    lines += ["## Short answers", "", "**Why no LLMs?** No credentials or local weights were available; simulation was chosen explicitly to produce reproducible causal-mechanism evidence without fabricating model results.", "", "**Why is H1 mixed?** Early faults had more reachable descendants, so EPR and depth rose. Late faults acted directly on the sink, so DAF was slightly higher despite zero downstream nodes.", "", "**Is full-trace attribution realistic?** No. It contains oracle validity and is an upper bound. A real replication must hide oracle labels and expose only realistic logs.", "", "**Why no composite score?** Weights encode stakeholder utility and can hide rank reversals. Component metrics are primary; any utility model must be declared separately.", "", "**Main contribution?** A unified, executable, auditable methodology connecting intervention integrity, propagation, attribution, mitigation, and reliability–cost analysis."]
    NOTES.write_text("\n".join(lines), encoding="utf-8")

if __name__ == "__main__":
    make_paper(); make_deck(); make_notes()
    print(PAPER); print(DECK); print(NOTES)
